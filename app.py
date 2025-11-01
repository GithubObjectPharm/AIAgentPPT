from flask import Flask, request, jsonify, send_from_directory
from werkzeug.utils import secure_filename
import os, re, shutil
from datetime import datetime
from docx import Document
from PyPDF2 import PdfReader
import openai

# Import pptx modules
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ---------------- CONFIG ----------------
app = Flask(__name__, static_folder="", template_folder="")

UPLOAD_FOLDER = "uploads"
GENERATED_FOLDER = "generated"
ALLOWED_EXTENSIONS = {"txt", "docx", "pdf"}

app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["GENERATED_FOLDER"] = GENERATED_FOLDER

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GENERATED_FOLDER, exist_ok=True)

openai.api_key = "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"

# ---------------- HELPERS ----------------
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text(filepath):
    """Extract readable text from txt, docx, or pdf"""
    ext = filepath.split(".")[-1].lower()
    try:
        if ext == "txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        elif ext == "docx":
            doc = Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])
        elif ext == "pdf":
            reader = PdfReader(filepath)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text
        else:
            return ""
    except Exception as e:
        print(f"Error extracting text: {e}")
        return ""

def extract_mcqs_simple(text):
    """Enhanced MCQ extraction that tolerates inconsistent formatting."""
    mcqs = []
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\n+', '\n', text)

    # Split by likely question indicators or bullet/question marks
    question_blocks = re.split(
        r'(?:\n\s*(?:Q\s*\d*[:\.\)]|\d+[\.\)]|\-?\s*Question\s*\d*|#|- |\• |\u2022 ))',
        text
    )

    for block in question_blocks:
        block = block.strip()
        if not block or len(block) < 10:
            continue

        lines = [line.strip() for line in block.split("\n") if line.strip()]
        if not lines:
            continue

        # Identify question (first non-option line)
        question = ""
        for ln in lines:
            if not re.match(r'^[A-Da-d][\.\)]|^[-•–]', ln):
                question = ln.strip()
                break
        if not question:
            continue

        # Collect options (A–D or bullet-style)
        options, answer, explanation = [], "", ""
        for line in lines:
            # Flexible option patterns
            opt_match = re.match(
                r'^[\-\s•–]*(?:([A-Da-d])[\.\)\:\-]\s*|[-•–]\s*)(.+)', line)
            if opt_match:
                label = opt_match.group(1).upper() if opt_match.group(1) else chr(65 + len(options))
                options.append(f"{label}) {opt_match.group(2).strip()}")
            elif re.search(r'(Answer|Ans|Correct|Option)\s*[:=\-]?\s*([A-Da-d])', line, re.IGNORECASE):
                answer = re.search(r'(Answer|Ans|Correct|Option)\s*[:=\-]?\s*([A-Da-d])', line, re.IGNORECASE).group(2).upper()
            elif re.search(r'(Explanation|Exp|Reason)\s*[:\-]?\s*(.+)', line, re.IGNORECASE):
                explanation = re.search(r'(Explanation|Exp|Reason)\s*[:\-]?\s*(.+)', line, re.IGNORECASE).group(2).strip()

        # Ensure we have usable defaults
        if not options:
            # Guess options from next few short lines
            guess_opts = [l for l in lines[1:5] if len(l.split()) <= 10]
            options = [f"{chr(65+i)}) {o}" for i, o in enumerate(guess_opts)] or ["A) Option 1", "B) Option 2"]

        mcqs.append({
            "question": question or "Question text missing",
            "options": options[:4],
            "answer": answer or "N/A",
            "explanation": explanation or "No explanation provided."
        })

    return mcqs


def create_vba_template_presentation(mcqs, output_path):
    """
    Create 100% exact replica of VBA template with improved spacing and NO animations
    (Animations removed to fix the error)
    """
    if not PPTX_AVAILABLE:
        return False
        
    try:
        # Create new presentation
        prs = Presentation()
        
        # Set slide size to match VBA template
        prs.slide_width = Inches(13.333)  # 16:9 widescreen
        prs.slide_height = Inches(7.5)
        
        # Exact colors from VBA code
        BLUE = RGBColor(100, 149, 237)      # Light blue borders
        CUSTOMBG = RGBColor(249, 248, 242)  # Background color
        DARK = RGBColor(22, 22, 26)         # Dark text
        WHITE = RGBColor(255, 255, 255)     # White text
        
        for i, mcq in enumerate(mcqs, 1):
            # ==================== SLIDE 1: QUESTION + OPTIONS ====================
            slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Background (exact color from VBA)
            bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg1.fill.solid()
            bg1.fill.fore_color.rgb = CUSTOMBG
            bg1.line.fill.background()
            
            # Question Box (rounded rectangle) - Larger and better positioned
            qBox1 = slide1.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(0.5), Inches(0.7), 
                Inches(12.3), Inches(2.0)  # Increased height for better spacing
            )
            qBox1.fill.solid()
            qBox1.fill.fore_color.rgb = CUSTOMBG
            qBox1.line.color.rgb = BLUE
            qBox1.line.width = Pt(8)  # Thicker border
            qBox1.adjustments[0] = 0.08  # Corner radius
            
            # Question Text
            qText1 = qBox1.text_frame
            qText1.text = f"Q{i}: {mcq['question']}"
            qText1.paragraphs[0].font.name = "Arial"
            qText1.paragraphs[0].font.size = Pt(36)
            qText1.paragraphs[0].font.bold = True
            qText1.paragraphs[0].font.color.rgb = DARK
            qText1.margin_left = Pt(25)
            qText1.margin_right = Pt(25)
            qText1.margin_top = Pt(20)
            qText1.margin_bottom = Pt(20)
            
            # Number Badge (oval) - Better positioned
            badge1 = slide1.shapes.add_shape(
                MSO_SHAPE.OVAL, 
                Inches(0.15), Inches(0.45), 
                Inches(0.9), Inches(0.9)  # Slightly larger badge
            )
            badge1.fill.solid()
            badge1.fill.fore_color.rgb = BLUE
            badge1.line.fill.background()
            
            # Badge Text
            badgeText1 = badge1.text_frame
            badgeText1.text = str(i)
            badgeText1.paragraphs[0].font.name = "Arial"
            badgeText1.paragraphs[0].font.size = Pt(36)  # Slightly larger
            badgeText1.paragraphs[0].font.bold = True
            badgeText1.paragraphs[0].font.color.rgb = WHITE
            badgeText1.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Answer Options with better spacing (animations removed)
            option_top = Inches(3.0)  # More space from question
            
            # Ensure we have exactly 4 options
            options = mcq['options']
            while len(options) < 4:
                options.append(f"{chr(65 + len(options))}) Option placeholder")
            
            for j, option in enumerate(options[:4]):  # Only take first 4 options
                option_box = slide1.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.5), option_top + (j * Inches(1.0)),  # Increased spacing between options
                    Inches(12.3), Inches(0.85)  # Taller option boxes
                )
                option_box.fill.solid()
                option_box.fill.fore_color.rgb = CUSTOMBG
                option_box.line.color.rgb = BLUE
                option_box.line.width = Pt(8)  # Thicker border
                option_box.adjustments[0] = 0.08
                
                option_text = option_box.text_frame
                option_text.text = option
                option_text.paragraphs[0].font.name = "Arial"
                option_text.paragraphs[0].font.size = Pt(28)
                option_text.paragraphs[0].font.italic = True
                option_text.paragraphs[0].font.color.rgb = DARK
                option_text.margin_left = Pt(25)
                option_text.margin_right = Pt(25)
                option_text.margin_top = Pt(15)
                option_text.margin_bottom = Pt(15)
            
            # ==================== SLIDE 2: ANSWER + EXPLANATION ====================
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Background
            bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg2.fill.solid()
            bg2.fill.fore_color.rgb = CUSTOMBG
            bg2.line.fill.background()
            
            # Question Box (same as slide 1)
            qBox2 = slide2.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(0.5), Inches(0.7), 
                Inches(12.3), Inches(2.0)
            )
            qBox2.fill.solid()
            qBox2.fill.fore_color.rgb = CUSTOMBG
            qBox2.line.color.rgb = BLUE
            qBox2.line.width = Pt(8)  # Thicker border
            qBox2.adjustments[0] = 0.08
            
            qText2 = qBox2.text_frame
            qText2.text = f"Q{i}: {mcq['question']}"
            qText2.paragraphs[0].font.name = "Arial"
            qText2.paragraphs[0].font.size = Pt(36)
            qText2.paragraphs[0].font.bold = True
            qText2.paragraphs[0].font.color.rgb = DARK
            qText2.margin_left = Pt(25)
            qText2.margin_right = Pt(25)
            qText2.margin_top = Pt(20)
            qText2.margin_bottom = Pt(20)
            
            # Number Badge (same as slide 1)
            badge2 = slide2.shapes.add_shape(
                MSO_SHAPE.OVAL, 
                Inches(0.15), Inches(0.45), 
                Inches(0.9), Inches(0.9)
            )
            badge2.fill.solid()
            badge2.fill.fore_color.rgb = BLUE
            badge2.line.fill.background()
            
            badgeText2 = badge2.text_frame
            badgeText2.text = str(i)
            badgeText2.paragraphs[0].font.name = "Arial"
            badgeText2.paragraphs[0].font.size = Pt(36)
            badgeText2.paragraphs[0].font.bold = True
            badgeText2.paragraphs[0].font.color.rgb = WHITE
            badgeText2.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Answer Explanation Box - Better spacing
            answer_box = slide2.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(3.0),  # Better positioned
                Inches(12.3), Inches(3.2)  # Better proportions
            )
            answer_box.fill.solid()
            answer_box.fill.fore_color.rgb = CUSTOMBG
            answer_box.line.color.rgb = BLUE
            answer_box.line.width = Pt(8)  # Thicker border
            answer_box.adjustments[0] = 0.08
            
            # Build answer content
            answer_content = []
            if mcq['answer']:
                answer_content.append(f"Correct Answer: {mcq['answer']}")
            if mcq['explanation']:
                answer_content.append(f"Explanation: {mcq['explanation']}")
            
            answer_text = answer_box.text_frame
            answer_text.text = "\n\n".join(answer_content)
            answer_text.margin_left = Pt(30)
            answer_text.margin_right = Pt(30)
            answer_text.margin_top = Pt(25)
            answer_text.margin_bottom = Pt(25)
            
            # Format answer text
            for paragraph in answer_text.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(28)
                paragraph.font.color.rgb = DARK
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Make first paragraph (answer) bold and italic like VBA
            if answer_text.paragraphs:
                answer_text.paragraphs[0].font.bold = True
                answer_text.paragraphs[0].font.italic = True
        
        # Save presentation
        prs.save(output_path)
        print(f"Successfully created VBA template presentation with {len(mcqs)} questions")
        return True
        
    except Exception as e:
        print(f"Error creating VBA template presentation: {e}")
        import traceback
        print(f"Detailed error: {traceback.format_exc()}")
        return False

def create_ppt_template_presentation(mcqs, output_path):
    """
    Create presentation using the existing PPT template file
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
        
        # Load the existing template
        prs = Presentation("templates/ppt_template.pptx")
        
        # Find a usable layout
        layout = None
        for slide_layout in prs.slide_layouts:
            if slide_layout.placeholders:
                layout = slide_layout
                break
        
        if not layout:
            layout = prs.slide_layouts[0]
        
        # Clear existing slides and create new ones
        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
        
        for i, mcq in enumerate(mcqs, 1):
            # Create question slide
            slide = prs.slides.add_slide(layout)
            
            # Set title
            if slide.shapes.title:
                title_text = f"Q{i}: {mcq['question']}"
                slide.shapes.title.text = title_text
            
            # Set content
            content_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Body/content placeholder
                    content_shape = shape
                    break
            
            if content_shape:
                tf = content_shape.text_frame
                tf.clear()
                
                # Add options
                for option in mcq['options']:
                    p = tf.add_paragraph()
                    p.text = option
                    p.font.size = Pt(18)
                
                # Add answer and explanation
                if mcq['answer'] or mcq['explanation']:
                    p = tf.add_paragraph()
                    answer_text = f"Answer: {mcq['answer']}" if mcq['answer'] else ""
                    explanation_text = f"Explanation: {mcq['explanation']}" if mcq['explanation'] else ""
                    p.text = f"{answer_text}\n{explanation_text}".strip()
                    p.font.bold = True
                    p.font.size = Pt(16)
        
        prs.save(output_path)
        return True
        
    except Exception as e:
        print(f"Error creating PPT template presentation: {e}")
        return False

# ---------------- ROUTES ----------------
@app.route("/")
def home():
    return send_from_directory(".", "index.html")

@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(filepath)
        return jsonify({"message": "✅ File uploaded successfully", "filename": filename})
    else:
        return jsonify({"error": "Invalid file type"}), 400

@app.route("/generate", methods=["POST"])
def generate_output():
    try:
        data = request.get_json()
        filename = data.get("filename")
        template_choice = data.get("template")

        if not filename or not template_choice:
            return jsonify({"error": "Missing parameters"}), 400

        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)

        if not os.path.exists(filepath):
            return jsonify({"error": f"Uploaded file not found: {filepath}"}), 404

        # Extract text content
        file_content = extract_text(filepath)
        if not file_content.strip():
            return jsonify({"error": "File is empty or unreadable."}), 400

        # Parse for MCQs
        mcqs = extract_mcqs_simple(file_content)

        # Fallback to AI if no MCQs found
        if not mcqs:
            try:
                gpt_response = openai.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "You are an educational assistant. Extract or create multiple choice questions from the provided content. Format each as: Q) question? A) option A B) option B C) option C D) option D Answer: [letter] Explanation: [text]"},
                        {"role": "user", "content": file_content[:8000]}
                    ]
                )
                gpt_output = gpt_response.choices[0].message.content
                mcqs = extract_mcqs_simple(gpt_output)
            except Exception as e:
                mcqs = [{
                    "question": "Content Summary",
                    "options": ["Key information extracted from uploaded file"],
                    "answer": "",
                    "explanation": file_content[:1000] + "..." if len(file_content) > 1000 else file_content
                }]

        # Generate output filename
        output_filename = f"generated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        output_path = os.path.join(app.config["GENERATED_FOLDER"], output_filename)

        # Create presentation based on template choice
        success = False
        if template_choice == "vba":
            success = create_vba_template_presentation(mcqs, output_path)
        elif template_choice == "ppt":
            success = create_ppt_template_presentation(mcqs, output_path)
        else:
            return jsonify({"error": "Invalid template choice"}), 400

        if success:
            return jsonify({
                "message": f"✅ {len(mcqs)} questions generated using {template_choice.upper()} template",
                "download_url": f"/download/{output_filename}",
                "questions_found": len(mcqs),
                "total_slides": len(mcqs) * 2 if template_choice == "vba" else len(mcqs)
            })
        else:
            return jsonify({"error": f"Failed to generate {template_choice} presentation"}), 500

    except Exception as e:
        return jsonify({"error": f"Unexpected error: {str(e)}"}), 500

@app.route("/download/<filename>")
def download_file(filename):
    return send_from_directory(app.config["GENERATED_FOLDER"], filename, as_attachment=True)

# ---------------- RUN APP ----------------
if __name__ == "__main__":
    if not PPTX_AVAILABLE:
        print("WARNING: python-pptx is not installed. Please run: pip install python-pptx")
    app.run(debug=True)